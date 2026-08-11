from __future__ import annotations

import asyncio
import io
import json
import unittest
from pathlib import Path

import httpx
from armi_kernel.application import (
    CredentialLocator,
    ExternalContentRecognitionRequest,
    ExternalContentRecognitionResult,
    ExternalContentRecognitionStatus,
    ExternalMessagePartKind,
    ExternalMessageViolation,
)
from armi_kernel.contracts import TraceId
from armi_runtime.adapters.model.doubao_speech import (
    DoubaoSpeechRecognizer,
    _audio_format,
)
from armi_runtime.adapters.model.external_content import (
    _input_message,
    load_external_recognition_binding,
)
from armi_runtime.composition.configuration import EnvironmentFileCredentialPort
from armi_runtime.composition.external_content_extractors import (
    extract_external_content,
)
from armi_runtime.composition.external_content_recognizer import (
    ExternalContentRecognizer,
)
from docx import Document
from openpyxl import Workbook
from pptx import Presentation


class ExternalContentExtractorTests(unittest.TestCase):
    def test_extracts_text_json_csv_and_gb18030(self) -> None:
        text = extract_external_content(
            kind=ExternalMessagePartKind.FILE,
            content="中文".encode("gb18030"),
            file_name="note.txt",
        )
        document = extract_external_content(
            kind=ExternalMessagePartKind.FILE,
            content=b'{"answer": 42}',
            file_name="value.json",
        )
        table = extract_external_content(
            kind=ExternalMessagePartKind.FILE,
            content=b"a,b\n1,2\n",
            file_name="value.csv",
        )
        self.assertEqual(text.text, "中文")
        self.assertIn('"answer": 42', document.text or "")
        self.assertEqual(table.text, "a | b\n1 | 2")

    def test_detects_office_container_independently_of_extension(self) -> None:
        docx_value = io.BytesIO()
        document = Document()
        document.add_paragraph("正文")
        document.save(docx_value)
        pptx_value = io.BytesIO()
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = "标题"
        presentation.save(pptx_value)
        xlsx_value = io.BytesIO()
        workbook = Workbook()
        workbook.active.title = "数据"
        workbook.active.append(("字段", "值"))
        workbook.save(xlsx_value)
        workbook.close()

        extracted = (
            extract_external_content(
                kind=ExternalMessagePartKind.FILE,
                content=value,
                file_name="opaque.bin",
            )
            for value in (
                docx_value.getvalue(),
                pptx_value.getvalue(),
                xlsx_value.getvalue(),
            )
        )
        docx, pptx, xlsx = extracted
        self.assertEqual(docx.text, "正文")
        self.assertIn("wordprocessingml", docx.media_type)
        self.assertIn("[幻灯片 1]\n标题", pptx.text or "")
        self.assertIn("[工作表 数据]\n字段 | 值", xlsx.text or "")

    def test_routes_pdf_and_rejects_mismatched_media(self) -> None:
        pdf = extract_external_content(
            kind=ExternalMessagePartKind.FILE,
            content=b"%PDF-1.7\nprobe",
            file_name="document.bin",
        )
        self.assertTrue(pdf.requires_provider)
        self.assertEqual(pdf.media_type, "application/pdf")
        with self.assertRaisesRegex(
            ExternalMessageViolation, "EXTERNAL-MESSAGE-MEDIA-TYPE"
        ):
            extract_external_content(
                kind=ExternalMessagePartKind.IMAGE,
                content=b"not-an-image",
                file_name="image.png",
            )


class ExternalContentModelRequestTests(unittest.TestCase):
    def test_loads_packaged_recognition_models(self) -> None:
        binding = load_external_recognition_binding(
            Path(
                "apps/armi-runtime/src/armi_runtime/composition/"
                "runtime_resources/model-bindings.manifest.json"
            )
        )
        self.assertEqual(
            binding.target_for(ExternalMessagePartKind.IMAGE),
            ("volcengine_ark", "doubao-seed-2-0-lite-260428"),
        )
        self.assertEqual(
            binding.target_for(ExternalMessagePartKind.AUDIO),
            ("volcengine_doubao_speech", "bigmodel-400"),
        )
        self.assertEqual(
            binding.target_for(ExternalMessagePartKind.VIDEO),
            ("volcengine_ark", "doubao-seed-2-0-lite-260428"),
        )
        self.assertEqual(
            binding.target_for(ExternalMessagePartKind.FILE),
            ("volcengine_ark", "doubao-seed-evolving"),
        )

    def test_builds_base64_requests_for_each_provider_media_kind(self) -> None:
        expected_types = {
            ExternalMessagePartKind.IMAGE: "input_image",
            ExternalMessagePartKind.VIDEO: "input_video",
            ExternalMessagePartKind.FILE: "input_file",
        }
        for kind, expected in expected_types.items():
            with self.subTest(kind=kind):
                value = _input_message(
                    ExternalContentRecognitionRequest(
                        kind,
                        b"sample",
                        "sample.bin",
                        "application/octet-stream",
                        TraceId("1" * 32),
                    )
                )
                self.assertEqual(value["content"][1]["type"], expected)
                self.assertNotIn("file_id", json.dumps(value))


class DoubaoSpeechRecognizerTests(unittest.TestCase):
    def _recognizer(self, handler) -> DoubaoSpeechRecognizer:
        bindings = load_external_recognition_binding(
            Path(
                "apps/armi-runtime/src/armi_runtime/composition/"
                "runtime_resources/model-bindings.manifest.json"
            )
        )
        return DoubaoSpeechRecognizer(
            credential_port=EnvironmentFileCredentialPort(
                environment={
                    "ARMI_SECRET_SPEECH_TEST": json.dumps(
                        {
                            "app_id": "test-speech-app",
                            "access_token": "test-speech-token",
                        }
                    )
                },
                secret_roots=(Path.cwd(),),
            ),
            locator=CredentialLocator.parse("env:ARMI_SECRET_SPEECH_TEST"),
            binding=bindings.speech,
            transport=httpx.MockTransport(handler),
        )

    def test_submits_standard_400_asr_and_returns_queried_transcript(self) -> None:
        observed: list[httpx.Request] = []

        def handler(request: httpx.Request) -> httpx.Response:
            observed.append(request)
            if request.url.path.endswith("/submit"):
                return httpx.Response(
                    200,
                    headers={
                        "X-Api-Status-Code": "20000000",
                        "X-Tt-Logid": "speech-submit-log-id",
                    },
                    json={},
                )
            return httpx.Response(
                200,
                headers={
                    "X-Api-Status-Code": "20000000",
                    "X-Tt-Logid": "speech-log-id",
                },
                json={
                    "audio_info": {"duration": 1200},
                    "result": {"text": "Hello, ARMI.", "utterances": []},
                },
            )

        result = asyncio.run(self._recognizer(handler).recognize(_audio_request()))

        self.assertEqual(result.status, ExternalContentRecognitionStatus.SUCCEEDED)
        self.assertEqual(result.text, "Hello, ARMI.")
        self.assertEqual(result.provider, "volcengine_doubao_speech")
        self.assertEqual(result.model_id, "bigmodel-400")
        self.assertEqual(result.provider_request_id, "speech-log-id")
        self.assertNotIn(b"test-speech-token", result.raw_response or b"")
        self.assertEqual(len(observed), 2)
        submitted, queried = observed
        self.assertTrue(submitted.url.path.endswith("/submit"))
        self.assertTrue(queried.url.path.endswith("/query"))
        self.assertEqual(submitted.headers["X-Api-App-Key"], "test-speech-app")
        self.assertEqual(submitted.headers["X-Api-Access-Key"], "test-speech-token")
        self.assertEqual(submitted.headers["X-Api-Resource-Id"], "volc.bigasr.auc")
        self.assertEqual(
            queried.headers["X-Api-Request-Id"],
            submitted.headers["X-Api-Request-Id"],
        )
        self.assertEqual(queried.headers["X-Tt-Logid"], "speech-submit-log-id")
        document = json.loads(submitted.content)
        self.assertEqual(document["user"]["uid"], "test-speech-app")
        self.assertEqual(document["audio"]["data"], "SUQz")
        self.assertEqual(document["audio"]["format"], "mp3")
        self.assertEqual(document["request"]["model_name"], "bigmodel")
        self.assertEqual(document["request"]["model_version"], "400")
        self.assertTrue(document["request"]["enable_itn"])
        self.assertTrue(document["request"]["enable_punc"])

    def test_maps_supported_audio_media_types_to_provider_formats(self) -> None:
        expected = {
            "audio/mpeg": "mp3",
            "audio/wav": "wav",
            "audio/ogg": "ogg",
            "audio/opus": "ogg",
        }
        for media_type, provider_format in expected.items():
            with self.subTest(media_type=media_type):
                request = ExternalContentRecognitionRequest(
                    ExternalMessagePartKind.AUDIO,
                    b"audio",
                    "sample",
                    media_type,
                    TraceId("1" * 32),
                )
                self.assertEqual(_audio_format(request), provider_format)

    def test_maps_provider_rejection_and_timeout_without_retry(self) -> None:
        calls = 0

        def rejected(_request: httpx.Request) -> httpx.Response:
            nonlocal calls
            calls += 1
            return httpx.Response(
                200,
                headers={"X-Api-Status-Code": "45000151"},
                json={"result": {}},
            )

        failed = asyncio.run(self._recognizer(rejected).recognize(_audio_request()))
        self.assertEqual(failed.status, ExternalContentRecognitionStatus.FAILED)
        self.assertEqual(failed.error_code, "EXTERNAL-MESSAGE-RECOGNITION-ASR-45000151")
        self.assertEqual(calls, 1)

        def timeout(request: httpx.Request) -> httpx.Response:
            raise httpx.ReadTimeout("timeout", request=request)

        unknown = asyncio.run(self._recognizer(timeout).recognize(_audio_request()))
        self.assertEqual(unknown.status, ExternalContentRecognitionStatus.UNKNOWN)

    def test_polls_queued_standard_task_without_resubmitting(self) -> None:
        calls: list[str] = []
        query_statuses = iter(("20000001", "20000002", "20000000"))

        def handler(request: httpx.Request) -> httpx.Response:
            calls.append(request.url.path)
            if request.url.path.endswith("/submit"):
                return httpx.Response(
                    200,
                    headers={"X-Api-Status-Code": "20000000"},
                    json={},
                )
            status = next(query_statuses)
            return httpx.Response(
                200,
                headers={"X-Api-Status-Code": status},
                json=(
                    {"result": {"text": "排队完成。"}} if status == "20000000" else {}
                ),
            )

        result = asyncio.run(self._recognizer(handler).recognize(_audio_request()))

        self.assertEqual(result.status, ExternalContentRecognitionStatus.SUCCEEDED)
        self.assertEqual(sum(path.endswith("/submit") for path in calls), 1)
        self.assertEqual(sum(path.endswith("/query") for path in calls), 3)

    def test_rejects_malformed_speech_credentials(self) -> None:
        bindings = load_external_recognition_binding(
            Path(
                "apps/armi-runtime/src/armi_runtime/composition/"
                "runtime_resources/model-bindings.manifest.json"
            )
        )
        recognizer = DoubaoSpeechRecognizer(
            credential_port=EnvironmentFileCredentialPort(
                environment={"ARMI_SECRET_SPEECH_TEST": "not-json"},
                secret_roots=(Path.cwd(),),
            ),
            locator=CredentialLocator.parse("env:ARMI_SECRET_SPEECH_TEST"),
            binding=bindings.speech,
            transport=httpx.MockTransport(
                lambda _request: self.fail("invalid credentials must not call ASR")
            ),
        )

        result = asyncio.run(recognizer.recognize(_audio_request()))

        self.assertEqual(result.status, ExternalContentRecognitionStatus.FAILED)
        self.assertEqual(result.error_code, "EXTERNAL-MESSAGE-RECOGNITION-CREDENTIAL")

    def test_router_uses_speech_only_for_audio(self) -> None:
        ark = _RecordingRecognizer("ark")
        speech = _RecordingRecognizer("speech")
        router = ExternalContentRecognizer(ark=ark, speech=speech)

        asyncio.run(router.recognize(_audio_request()))
        asyncio.run(
            router.recognize(
                ExternalContentRecognitionRequest(
                    ExternalMessagePartKind.VIDEO,
                    b"video",
                    "sample.mp4",
                    "video/mp4",
                    TraceId("2" * 32),
                )
            )
        )

        self.assertEqual(speech.kinds, [ExternalMessagePartKind.AUDIO])
        self.assertEqual(ark.kinds, [ExternalMessagePartKind.VIDEO])


class _RecordingRecognizer:
    def __init__(self, provider: str) -> None:
        self.provider = provider
        self.kinds: list[ExternalMessagePartKind] = []

    async def recognize(
        self, request: ExternalContentRecognitionRequest
    ) -> ExternalContentRecognitionResult:
        self.kinds.append(request.kind)
        return ExternalContentRecognitionResult(
            ExternalContentRecognitionStatus.SUCCEEDED,
            "ok",
            self.provider,
            "test-model",
            None,
            None,
            None,
            None,
            b"{}\n",
            None,
        )


def _audio_request() -> ExternalContentRecognitionRequest:
    return ExternalContentRecognitionRequest(
        ExternalMessagePartKind.AUDIO,
        b"ID3",
        "sample.mp3",
        "audio/mpeg",
        TraceId("1" * 32),
    )


if __name__ == "__main__":
    unittest.main()
