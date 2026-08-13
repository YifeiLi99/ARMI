from __future__ import annotations

import io
import struct
import unittest
import zlib

from armi_interaction.api import (
    ExternalMessagePartKind,
    ExternalMessageViolation,
    ExternalVisualRole,
)
from armi_kernel.contracts import TraceId
from armi_perception._extractors import extract_external_content
from armi_perception.api import (
    ExternalContentRecognitionRequest,
    ExternalContentRecognitionResult,
    ExternalContentRecognitionStatus,
)
from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation


class ExternalContentExtractorTests(unittest.TestCase):
    def test_detects_supported_static_formats_and_ignores_claimed_extension(
        self,
    ) -> None:
        expected = {
            "JPEG": "image/jpeg",
            "PNG": "image/png",
            "GIF": "image/gif",
            "WEBP": "image/webp",
            "BMP": "image/bmp",
        }
        for image_format, media_type in expected.items():
            with self.subTest(image_format=image_format):
                value = io.BytesIO()
                Image.new("RGB", (12, 10), "blue").save(value, format=image_format)
                extracted = extract_external_content(
                    kind=ExternalMessagePartKind.IMAGE,
                    content=value.getvalue(),
                    file_name="falsely-declared.jpg",
                    visual_role=ExternalVisualRole.ORDINARY,
                )
                self.assertEqual(extracted.media_type, media_type)
                self.assertEqual(
                    (
                        extracted.pixel_width,
                        extracted.pixel_height,
                        extracted.frame_count,
                    ),
                    (12, 10, 1),
                )

    def test_inspects_static_image_and_uses_meaningful_sticker_summary(self) -> None:
        value = io.BytesIO()
        Image.new("RGB", (76, 64), "red").save(value, format="JPEG")
        extracted = extract_external_content(
            kind=ExternalMessagePartKind.IMAGE,
            content=value.getvalue(),
            file_name="sticker.jpg",
            visual_role=ExternalVisualRole.STICKER,
            source_summary="开心企鹅",
        )
        self.assertEqual(
            (extracted.media_type, extracted.pixel_width, extracted.pixel_height),
            ("image/jpeg", 76, 64),
        )
        self.assertEqual(extracted.frame_count, 1)
        self.assertFalse(extracted.requires_provider)
        self.assertIn("开心企鹅", extracted.text or "")
        self.assertEqual(len(extracted.visual_inputs), 1)

        generic = extract_external_content(
            kind=ExternalMessagePartKind.IMAGE,
            content=value.getvalue(),
            file_name="sticker.jpg",
            visual_role=ExternalVisualRole.STICKER,
            source_summary="商城表情",
        )
        self.assertTrue(generic.requires_provider)

    def test_extracts_four_ordered_png_frames_from_animated_gif(self) -> None:
        value = io.BytesIO()
        frames = [
            Image.new("RGBA", (24, 24), (index * 30, 0, 0, 255)) for index in range(6)
        ]
        frames[0].save(
            value,
            format="GIF",
            save_all=True,
            append_images=frames[1:],
            duration=20,
            loop=0,
        )
        extracted = extract_external_content(
            kind=ExternalMessagePartKind.IMAGE,
            content=value.getvalue(),
            file_name="animated.gif",
            visual_role=ExternalVisualRole.STICKER_CANDIDATE,
        )
        self.assertTrue(extracted.requires_provider)
        self.assertEqual(extracted.frame_count, 6)
        self.assertEqual(len(extracted.visual_inputs), 4)
        self.assertTrue(
            all(item.media_type == "image/png" for item in extracted.visual_inputs)
        )

    def test_extracts_text_json_csv_and_gb18030(self) -> None:
        text = extract_external_content(
            kind=ExternalMessagePartKind.FILE,
            content="中文".encode("gb18030"),
            file_name="note.txt",
        )
        document = extract_external_content(
            kind=ExternalMessagePartKind.FILE,
            content=b'{"answer": 42}',
            file_name="value.json",
        )
        table = extract_external_content(
            kind=ExternalMessagePartKind.FILE,
            content=b"a,b\n1,2\n",
            file_name="value.csv",
        )
        self.assertEqual(text.text, "中文")
        self.assertIn('"answer": 42', document.text or "")
        self.assertEqual(table.text, "a | b\n1 | 2")

    def test_detects_office_container_independently_of_extension(self) -> None:
        docx_value = io.BytesIO()
        document = Document()
        document.add_paragraph("正文")
        document.save(docx_value)
        pptx_value = io.BytesIO()
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title = slide.shapes.title
        assert title is not None
        title.text = "标题"
        presentation.save(pptx_value)
        xlsx_value = io.BytesIO()
        workbook = Workbook()
        worksheet = workbook.active
        assert worksheet is not None
        worksheet.title = "数据"
        worksheet.append(("字段", "值"))
        workbook.save(xlsx_value)
        workbook.close()

        extracted = (
            extract_external_content(
                kind=ExternalMessagePartKind.FILE,
                content=value,
                file_name="opaque.bin",
            )
            for value in (
                docx_value.getvalue(),
                pptx_value.getvalue(),
                xlsx_value.getvalue(),
            )
        )
        docx, pptx, xlsx = extracted
        self.assertEqual(docx.text, "正文")
        self.assertIn("wordprocessingml", docx.media_type)
        self.assertIn("[幻灯片 1]\n标题", pptx.text or "")
        self.assertIn("[工作表 数据]\n字段 | 值", xlsx.text or "")

    def test_routes_pdf_and_rejects_mismatched_media(self) -> None:
        pdf = extract_external_content(
            kind=ExternalMessagePartKind.FILE,
            content=b"%PDF-1.7\nprobe",
            file_name="document.bin",
        )
        self.assertTrue(pdf.requires_provider)
        self.assertEqual(pdf.media_type, "application/pdf")
        with self.assertRaisesRegex(
            ExternalMessageViolation, "EXTERNAL-MESSAGE-MEDIA-TYPE"
        ):
            extract_external_content(
                kind=ExternalMessagePartKind.IMAGE,
                content=b"not-an-image",
                file_name="image.png",
                visual_role=ExternalVisualRole.UNKNOWN,
            )

    def test_rejects_corrupt_and_over_36_megapixel_images(self) -> None:
        with self.assertRaisesRegex(
            ExternalMessageViolation, "EXTERNAL-MESSAGE-FILE-INVALID"
        ):
            extract_external_content(
                kind=ExternalMessagePartKind.IMAGE,
                content=b"\x89PNG\r\n\x1a\ncorrupt",
                file_name="corrupt.png",
                visual_role=ExternalVisualRole.UNKNOWN,
            )

        def chunk(kind: bytes, value: bytes) -> bytes:
            return (
                struct.pack(">I", len(value))
                + kind
                + value
                + struct.pack(">I", zlib.crc32(kind + value) & 0xFFFFFFFF)
            )

        oversized = (
            b"\x89PNG\r\n\x1a\n"
            + chunk(b"IHDR", struct.pack(">IIBBBBB", 6001, 6000, 8, 2, 0, 0, 0))
            + chunk(b"IEND", b"")
        )
        with self.assertRaisesRegex(
            ExternalMessageViolation, "EXTERNAL-MESSAGE-IMAGE-DIMENSIONS"
        ):
            extract_external_content(
                kind=ExternalMessagePartKind.IMAGE,
                content=oversized,
                file_name="oversized.png",
                visual_role=ExternalVisualRole.ORDINARY,
            )


class _RecordingRecognizer:
    def __init__(self, provider: str) -> None:
        self.provider = provider
        self.kinds: list[ExternalMessagePartKind] = []

    async def recognize(
        self, request: ExternalContentRecognitionRequest
    ) -> ExternalContentRecognitionResult:
        self.kinds.append(request.kind)
        return ExternalContentRecognitionResult(
            ExternalContentRecognitionStatus.SUCCEEDED,
            "ok",
            self.provider,
            "test-model",
            None,
            None,
            None,
            None,
            b"{}\n",
            None,
        )


def _audio_request() -> ExternalContentRecognitionRequest:
    return ExternalContentRecognitionRequest(
        ExternalMessagePartKind.AUDIO,
        b"ID3",
        "sample.mp3",
        "audio/mpeg",
        TraceId("1" * 32),
    )


if __name__ == "__main__":
    unittest.main()
