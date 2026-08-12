"""Deterministic extraction and media validation for external message parts."""

from __future__ import annotations

import csv
import io
import json
import warnings
from dataclasses import dataclass
from pathlib import Path
from zipfile import BadZipFile, ZipFile

from armi_interaction.api import (
    ExternalMediaContent,
    ExternalMessagePartKind,
    ExternalMessageViolation,
    ExternalVisualRole,
)
from docx import Document
from docx.opc.exceptions import PackageNotFoundError as DocxPackageNotFoundError
from openpyxl import load_workbook
from openpyxl.utils.exceptions import InvalidFileException
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError

_MAX_PROJECTION_BYTES = 256 * 1024
_TRUNCATED = "\n[内容已按 ARMI 单条认知材料上限截断]"
_MAX_IMAGE_PIXELS = 36_000_000
_GENERIC_STICKER_SUMMARIES = frozenset(
    {
        "图片",
        "表情",
        "表情包",
        "动画表情",
        "商城表情",
        "[图片]",
        "[表情]",
        "[动画表情]",
        "[商城表情]",
    }
)


@dataclass(frozen=True, slots=True)
class ExtractedExternalContent:
    media_type: str
    text: str | None
    requires_provider: bool
    pixel_width: int | None = None
    pixel_height: int | None = None
    frame_count: int | None = None
    visual_inputs: tuple[ExternalMediaContent, ...] = ()


def extract_external_content(
    *,
    kind: ExternalMessagePartKind,
    content: bytes,
    file_name: str,
    visual_role: ExternalVisualRole | None = None,
    source_summary: str | None = None,
) -> ExtractedExternalContent:
    if kind is ExternalMessagePartKind.IMAGE:
        if type(visual_role) is not ExternalVisualRole:
            raise ExternalMessageViolation("EXTERNAL-MESSAGE-MEDIA-TYPE")
        return _extract_image(
            content,
            file_name=file_name,
            visual_role=visual_role,
            source_summary=source_summary,
        )
    if kind is ExternalMessagePartKind.AUDIO:
        if not (content.startswith(b"ID3") or _looks_like_mp3_frame(content)):
            raise ExternalMessageViolation("EXTERNAL-MESSAGE-MEDIA-TYPE")
        return ExtractedExternalContent("audio/mpeg", None, True)
    if kind is ExternalMessagePartKind.VIDEO:
        if len(content) < 12 or content[4:8] != b"ftyp":
            raise ExternalMessageViolation("EXTERNAL-MESSAGE-MEDIA-TYPE")
        return ExtractedExternalContent("video/mp4", None, True)
    if kind is not ExternalMessagePartKind.FILE:
        raise ExternalMessageViolation("EXTERNAL-MESSAGE-MEDIA-TYPE")
    suffix = Path(file_name).suffix.lower()
    if content.startswith(b"%PDF-"):
        return ExtractedExternalContent("application/pdf", None, True)
    if suffix in {".txt", ".md", ".markdown", ".log", ".json", ".csv"} or suffix in {
        ".py",
        ".js",
        ".ts",
        ".tsx",
        ".jsx",
        ".java",
        ".c",
        ".h",
        ".cpp",
        ".cs",
        ".go",
        ".rs",
        ".sql",
        ".toml",
        ".yaml",
        ".yml",
        ".xml",
        ".html",
        ".css",
    }:
        text = _decode_text(content)
        if suffix == ".json":
            try:
                text = json.dumps(json.loads(text), ensure_ascii=False, indent=2)
            except json.JSONDecodeError:
                raise ExternalMessageViolation(
                    "EXTERNAL-MESSAGE-FILE-INVALID"
                ) from None
        elif suffix == ".csv":
            text = "\n".join(" | ".join(row) for row in csv.reader(io.StringIO(text)))
        return ExtractedExternalContent("text/plain", _bounded(text), False)
    try:
        office_kind = _office_kind(content)
        text = (
            _docx_text(content)
            if office_kind == "docx"
            else _pptx_text(content)
            if office_kind == "pptx"
            else _xlsx_text(content)
            if office_kind == "xlsx"
            else None
        )
    except (
        BadZipFile,
        DocxPackageNotFoundError,
        InvalidFileException,
        KeyError,
        OSError,
        PptxPackageNotFoundError,
        SyntaxError,
        ValueError,
    ):
        raise ExternalMessageViolation("EXTERNAL-MESSAGE-FILE-INVALID") from None
    if office_kind == "docx":
        return ExtractedExternalContent(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            _bounded(text or ""),
            False,
        )
    if office_kind == "pptx":
        return ExtractedExternalContent(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            _bounded(text or ""),
            False,
        )
    if office_kind == "xlsx":
        return ExtractedExternalContent(
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            _bounded(text or ""),
            False,
        )
    return ExtractedExternalContent(
        "application/octet-stream",
        _bounded(f"[文件: {file_name},该文件类型暂不支持读取内容]"),
        False,
    )


def _office_kind(content: bytes) -> str | None:
    if not content.startswith(b"PK"):
        return None
    with ZipFile(io.BytesIO(content)) as archive:
        names = frozenset(archive.namelist())
    if "word/document.xml" in names:
        return "docx"
    if "ppt/presentation.xml" in names:
        return "pptx"
    if "xl/workbook.xml" in names:
        return "xlsx"
    return None


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8-sig", "gb18030"):
        try:
            return content.decode(encoding, errors="strict")
        except UnicodeDecodeError:
            continue
    raise ExternalMessageViolation("EXTERNAL-MESSAGE-FILE-INVALID")


def _bounded(value: str) -> str:
    if not value.strip():
        raise ExternalMessageViolation("EXTERNAL-MESSAGE-FILE-EMPTY")
    encoded = value.encode("utf-8", errors="strict")
    if len(encoded) <= _MAX_PROJECTION_BYTES:
        return value
    limit = _MAX_PROJECTION_BYTES - len(_TRUNCATED.encode("utf-8"))
    clipped = encoded[:limit]
    while True:
        try:
            return clipped.decode("utf-8", errors="strict") + _TRUNCATED
        except UnicodeDecodeError:
            clipped = clipped[:-1]


def _image_media_type(content: bytes) -> str:
    signatures = (
        (b"\xff\xd8\xff", "image/jpeg"),
        (b"\x89PNG\r\n\x1a\n", "image/png"),
        (b"GIF87a", "image/gif"),
        (b"GIF89a", "image/gif"),
        (b"BM", "image/bmp"),
    )
    for signature, media_type in signatures:
        if content.startswith(signature):
            return media_type
    if len(content) >= 12 and content.startswith(b"RIFF") and content[8:12] == b"WEBP":
        return "image/webp"
    raise ExternalMessageViolation("EXTERNAL-MESSAGE-MEDIA-TYPE")


def _extract_image(
    content: bytes,
    *,
    file_name: str,
    visual_role: ExternalVisualRole,
    source_summary: str | None,
) -> ExtractedExternalContent:
    media_type = _image_media_type(content)
    try:
        with warnings.catch_warnings():
            warnings.simplefilter("error", Image.DecompressionBombWarning)
            with Image.open(io.BytesIO(content)) as image:
                width, height = image.size
                frame_count = int(getattr(image, "n_frames", 1))
                if (
                    width <= 0
                    or height <= 0
                    or width * height > _MAX_IMAGE_PIXELS
                    or frame_count <= 0
                ):
                    raise ExternalMessageViolation("EXTERNAL-MESSAGE-IMAGE-DIMENSIONS")
                image.verify()
            inputs = _visual_inputs(
                content,
                file_name=file_name,
                media_type=media_type,
                frame_count=frame_count,
            )
    except ExternalMessageViolation:
        raise
    except Image.DecompressionBombError, Image.DecompressionBombWarning:
        raise ExternalMessageViolation("EXTERNAL-MESSAGE-IMAGE-DIMENSIONS") from None
    except OSError, SyntaxError, UnidentifiedImageError, ValueError:
        raise ExternalMessageViolation("EXTERNAL-MESSAGE-FILE-INVALID") from None
    local_text = None
    if visual_role is ExternalVisualRole.STICKER and _meaningful_sticker_summary(
        source_summary
    ):
        local_text = f"QQ 提供的商城表情摘要: {source_summary}"
    return ExtractedExternalContent(
        media_type,
        local_text,
        local_text is None,
        width,
        height,
        frame_count,
        inputs,
    )


def _visual_inputs(
    content: bytes, *, file_name: str, media_type: str, frame_count: int
) -> tuple[ExternalMediaContent, ...]:
    if frame_count == 1 and media_type in {"image/jpeg", "image/png"}:
        return (ExternalMediaContent(content, file_name, media_type),)
    indexes = _frame_indexes(frame_count)
    values: list[ExternalMediaContent] = []
    with Image.open(io.BytesIO(content)) as image:
        for ordinal, index in enumerate(indexes, start=1):
            image.seek(index)
            frame = image.convert("RGBA")
            output = io.BytesIO()
            frame.save(output, format="PNG")
            values.append(
                ExternalMediaContent(
                    output.getvalue(),
                    f"{Path(file_name).stem}-frame-{ordinal}.png",
                    "image/png",
                )
            )
    return tuple(values)


def _frame_indexes(frame_count: int) -> tuple[int, ...]:
    if frame_count <= 4:
        return tuple(range(frame_count))
    return tuple(round(index * (frame_count - 1) / 3) for index in range(4))


def _meaningful_sticker_summary(value: str | None) -> bool:
    if value is None:
        return False
    normalized = value.strip()
    return bool(normalized) and normalized.casefold() not in {
        item.casefold() for item in _GENERIC_STICKER_SUMMARIES
    }


def _looks_like_mp3_frame(content: bytes) -> bool:
    return len(content) >= 2 and content[0] == 0xFF and content[1] & 0xE0 == 0xE0


def _docx_text(content: bytes) -> str:
    document = Document(io.BytesIO(content))
    chunks = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    for table_index, table in enumerate(document.tables, start=1):
        chunks.append(f"[表格 {table_index}]")
        chunks.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)
    return "\n".join(chunks)


def _pptx_text(content: bytes) -> str:
    presentation = Presentation(io.BytesIO(content))
    chunks: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        values: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str) and text:
                values.append(text)
        chunks.append(f"[幻灯片 {index}]\n" + "\n".join(values))
    return "\n\n".join(chunks)


def _xlsx_text(content: bytes) -> str:
    workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=False)
    chunks: list[str] = []
    try:
        for sheet in workbook.worksheets:
            chunks.append(f"[工作表 {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                chunks.append(
                    " | ".join("" if value is None else str(value) for value in row)
                )
    finally:
        workbook.close()
    return "\n".join(chunks)


__all__ = ("ExtractedExternalContent", "extract_external_content")
